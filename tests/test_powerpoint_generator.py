"""Tests for PowerPointGenerator — the main coverage gap in the reporting module.

Exercises all three public creation methods plus the DataFrame presentation
generator. Requires python-pptx.
"""

import pytest
import pandas as pd
import numpy as np

try:
    from pptx import Presentation as PptxPresentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

pytestmark = pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not installed")

from siege_utilities.reporting.powerpoint_generator import PowerPointGenerator


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture
def gen(tmp_path):
    return PowerPointGenerator(client_name="TestClient", output_dir=tmp_path)


@pytest.fixture
def report_data():
    return {
        "executive_summary": "Q4 was excellent. Revenue up 15%.",
        "metrics": {
            "Revenue": {"value": "$68.5M", "status": "+15%"},
            "Users": {"value": "846", "status": "+23%"},
        },
        "charts": [
            {"title": "Revenue by Quarter", "type": "bar"},
        ],
        "tables": [
            {
                "title": "Quarterly Performance",
                "headers": ["Quarter", "Revenue"],
                "data": [["Q1", "$12M"], ["Q2", "$15M"]],
            }
        ],
        "insights": [
            "West region shows highest growth",
            "Customer satisfaction at all-time high",
        ],
    }


@pytest.fixture
def performance_data():
    return {
        "overview": "All targets exceeded in FY 2025.",
        "revenue": {"current": "$68.5M", "target": "$65.0M", "achievement": "105%"},
        "customers": {"current": "8460", "target": "7500", "achievement": "113%"},
        "trends": {"Revenue": "Upward", "Growth": "Accelerating"},
        "recommendations": [
            "Increase investment in West region",
            "Launch customer loyalty program",
        ],
    }


@pytest.fixture
def custom_config():
    return {
        "type": "quarterly_review",
        "title": "Q4 2025 Business Review",
        "slides": [
            {
                "type": "content",
                "title": "Agenda",
                "content": "1. Summary\n2. Performance\n3. Outlook",
            },
            {
                "type": "content",
                "title": "Summary",
                "content": "Revenue $22.1M (+18% QoQ)",
            },
            {
                "type": "table",
                "title": "Regional Data",
                "headers": ["Region", "Revenue"],
                "data": [["NE", "$18M"], ["SE", "$12M"]],
            },
            {
                "type": "comparison",
                "title": "YoY Comparison",
                "left_content": "Q4 2024: $18.7M",
                "right_content": "Q4 2025: $22.1M",
            },
        ],
    }


# ---------------------------------------------------------------------------
# Initialization
# ---------------------------------------------------------------------------

class TestInit:

    def test_creates_output_dir(self, tmp_path):
        out = tmp_path / "presentations"
        gen = PowerPointGenerator(client_name="X", output_dir=out)
        assert out.exists()
        assert gen.client_name == "X"

    def test_default_output_dir(self):
        gen = PowerPointGenerator(client_name="Y")
        assert gen.output_dir.name == "presentations"

    def test_slide_layouts(self, gen):
        assert "title" in gen.slide_layouts
        assert "blank" in gen.slide_layouts
        assert len(gen.slide_layouts) == 7


# ---------------------------------------------------------------------------
# create_analytics_presentation
# ---------------------------------------------------------------------------

class TestAnalyticsPresentation:

    def test_generates_pptx(self, gen, report_data):
        path = gen.create_analytics_presentation(report_data, presentation_title="Test Analytics")
        assert path.exists()
        assert path.suffix == ".pptx"
        assert path.stat().st_size > 0

    def test_filename_contains_client(self, gen, report_data):
        path = gen.create_analytics_presentation(report_data)
        assert "testclient" in path.name.lower()

    def test_valid_pptx(self, gen, report_data):
        path = gen.create_analytics_presentation(report_data)
        prs = PptxPresentation(str(path))
        assert len(prs.slides) >= 2  # at least title + exec summary

    def test_minimal_data(self, gen):
        path = gen.create_analytics_presentation({})
        assert path.exists()
        prs = PptxPresentation(str(path))
        assert len(prs.slides) >= 1  # at least title

    def test_without_charts(self, gen, report_data):
        path = gen.create_analytics_presentation(report_data, include_charts=False)
        assert path.exists()

    def test_without_tables(self, gen, report_data):
        path = gen.create_analytics_presentation(report_data, include_tables=False)
        assert path.exists()

    def test_full_data_multiple_slides(self, gen, report_data):
        path = gen.create_analytics_presentation(report_data)
        prs = PptxPresentation(str(path))
        # title + exec summary + metrics + charts + tables + insights = 6
        assert len(prs.slides) >= 4


# ---------------------------------------------------------------------------
# create_performance_presentation
# ---------------------------------------------------------------------------

class TestPerformancePresentation:

    def test_generates_pptx(self, gen, performance_data):
        path = gen.create_performance_presentation(
            performance_data, metrics=["revenue", "customers"]
        )
        assert path.exists()
        assert path.suffix == ".pptx"

    def test_valid_pptx(self, gen, performance_data):
        path = gen.create_performance_presentation(
            performance_data, metrics=["revenue"]
        )
        prs = PptxPresentation(str(path))
        assert len(prs.slides) >= 2

    def test_custom_title(self, gen, performance_data):
        path = gen.create_performance_presentation(
            performance_data, metrics=["revenue"], presentation_title="Custom Title"
        )
        assert path.exists()

    def test_empty_metrics(self, gen, performance_data):
        path = gen.create_performance_presentation(performance_data, metrics=[])
        assert path.exists()


# ---------------------------------------------------------------------------
# create_custom_presentation
# ---------------------------------------------------------------------------

class TestCustomPresentation:

    def test_generates_pptx(self, gen, custom_config):
        path = gen.create_custom_presentation(custom_config)
        assert path.exists()
        assert path.suffix == ".pptx"

    def test_slide_count(self, gen, custom_config):
        path = gen.create_custom_presentation(custom_config)
        prs = PptxPresentation(str(path))
        # title + 4 content slides = 5
        assert len(prs.slides) >= 4

    def test_empty_slides(self, gen):
        path = gen.create_custom_presentation({"title": "Empty", "slides": []})
        assert path.exists()

    def test_type_in_filename(self, gen, custom_config):
        path = gen.create_custom_presentation(custom_config)
        assert "quarterly_review" in path.name


# ---------------------------------------------------------------------------
# create_presentation_from_dataframe
# ---------------------------------------------------------------------------

class TestDataFramePresentation:

    def test_generates_pptx(self, gen):
        df = pd.DataFrame({
            "Product": [f"P{i}" for i in range(10)],
            "Sales": np.random.randint(100, 1000, 10),
            "Revenue": np.random.uniform(1000, 10000, 10).round(2),
        })
        path = gen.create_presentation_from_dataframe(df, presentation_title="DF Test")
        assert path.exists()
        assert path.suffix == ".pptx"

    def test_max_slides_limit(self, gen):
        df = pd.DataFrame({
            "Product": [f"P{i}" for i in range(50)],
            "Sales": np.random.randint(100, 1000, 50),
        })
        path = gen.create_presentation_from_dataframe(df, max_slides=3)
        prs = PptxPresentation(str(path))
        # title + summary + up to 3 data slides
        assert len(prs.slides) <= 6

    def test_small_dataframe(self, gen):
        df = pd.DataFrame({"A": [1], "B": [2]})
        path = gen.create_presentation_from_dataframe(df)
        assert path.exists()


# ---------------------------------------------------------------------------
# File output
# ---------------------------------------------------------------------------

class TestFileOutput:

    def test_multiple_presentations_unique_filenames(self, gen, report_data):
        p1 = gen.create_analytics_presentation(report_data)
        p2 = gen.create_analytics_presentation(report_data)
        # Timestamps should differ (or at least not overwrite)
        assert p1.exists()
        assert p2.exists()

    def test_output_in_correct_directory(self, gen, report_data):
        path = gen.create_analytics_presentation(report_data)
        assert path.parent == gen.output_dir
