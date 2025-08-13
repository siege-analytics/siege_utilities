"""
PowerPoint generation utilities for siege_utilities reporting system.
Creates presentations from analytics data and report configurations.
"""

import logging
from pathlib import Path
from typing import Dict, List, Any, Optional, Union
from datetime import datetime
import pandas as pd

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    Presentation = None

log = logging.getLogger(__name__)

class PowerPointGenerator:
    """
    Generates PowerPoint presentations from analytics data and report configurations.
    Requires python-pptx package for PowerPoint generation.
    """

    def __init__(self, client_name: str, output_dir: Optional[Path] = None):
        """
        Initialize the PowerPoint generator.
        
        Args:
            client_name: Name of the client for branding
            output_dir: Directory for output presentations
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx package is required for PowerPoint generation. Install with: pip install python-pptx")
        
        self.client_name = client_name
        self.output_dir = output_dir or Path.cwd() / "presentations"
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Default slide layouts
        self.slide_layouts = {
            'title': 0,
            'content': 1,
            'section_header': 2,
            'two_content': 3,
            'comparison': 4,
            'title_only': 5,
            'blank': 6
        }

    def create_analytics_presentation(self, report_data: Dict[str, Any],
                                   presentation_title: str = "",
                                   include_charts: bool = True,
                                   include_tables: bool = True) -> Path:
        """
        Create a PowerPoint presentation from analytics data.
        
        Args:
            report_data: Dictionary containing report data
            presentation_title: Title of the presentation
            include_charts: Whether to include charts
            include_tables: Whether to include data tables
            
        Returns:
            Path to the generated PowerPoint file
        """
        try:
            # Generate filename
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"{self.client_name.lower().replace(' ', '_')}_analytics_presentation_{timestamp}.pptx"
            output_path = self.output_dir / filename
            
            # Create presentation
            prs = Presentation()
            
            # Add title slide
            self._add_title_slide(prs, presentation_title or f"{self.client_name} Analytics Report")
            
            # Add executive summary slide
            if 'executive_summary' in report_data:
                self._add_executive_summary_slide(prs, report_data['executive_summary'])
            
            # Add metrics slide
            if 'metrics' in report_data:
                self._add_metrics_slide(prs, report_data['metrics'])
            
            # Add charts slides
            if include_charts and 'charts' in report_data:
                self._add_charts_slides(prs, report_data['charts'])
            
            # Add tables slides
            if include_tables and 'tables' in report_data:
                self._add_tables_slides(prs, report_data['tables'])
            
            # Add insights slide
            if 'insights' in report_data:
                self._add_insights_slide(prs, report_data['insights'])
            
            # Save presentation
            prs.save(str(output_path))
            
            log.info(f"PowerPoint presentation created: {output_path}")
            return output_path
            
        except Exception as e:
            log.error(f"Error creating PowerPoint presentation: {e}")
            raise

    def create_performance_presentation(self, performance_data: Dict[str, Any],
                                     metrics: List[str],
                                     presentation_title: str = "") -> Path:
        """
        Create a performance metrics presentation.
        
        Args:
            performance_data: Dictionary containing performance metrics
            metrics: List of metrics to include
            presentation_title: Custom title for the presentation
            
        Returns:
            Path to the generated PowerPoint file
        """
        try:
            # Generate filename
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"{self.client_name.lower().replace(' ', '_')}_performance_presentation_{timestamp}.pptx"
            output_path = self.output_dir / filename
            
            # Create presentation
            prs = Presentation()
            
            # Add title slide
            title = presentation_title or f"{self.client_name} Performance Report"
            self._add_title_slide(prs, title)
            
            # Add performance overview slide
            if 'overview' in performance_data:
                self._add_performance_overview_slide(prs, performance_data['overview'])
            
            # Add metrics breakdown slide
            self._add_metrics_breakdown_slide(prs, performance_data, metrics)
            
            # Add trends slides
            if 'trends' in performance_data:
                self._add_trends_slides(prs, performance_data['trends'])
            
            # Add recommendations slide
            if 'recommendations' in performance_data:
                self._add_recommendations_slide(prs, performance_data['recommendations'])
            
            # Save presentation
            prs.save(str(output_path))
            
            log.info(f"Performance PowerPoint presentation created: {output_path}")
            return output_path
            
        except Exception as e:
            log.error(f"Error creating performance PowerPoint presentation: {e}")
            raise

    def create_custom_presentation(self, presentation_config: Dict[str, Any]) -> Path:
        """
        Create a custom PowerPoint presentation based on configuration.
        
        Args:
            presentation_config: Complete presentation configuration
            
        Returns:
            Path to the generated PowerPoint file
        """
        try:
            # Generate filename
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            presentation_type = presentation_config.get('type', 'custom')
            filename = f"{self.client_name.lower().replace(' ', '_')}_{presentation_type}_presentation_{timestamp}.pptx"
            output_path = self.output_dir / filename
            
            # Create presentation
            prs = Presentation()
            
            # Add title slide
            title = presentation_config.get('title', f"{self.client_name} Presentation")
            self._add_title_slide(prs, title)
            
            # Add custom slides
            slides = presentation_config.get('slides', [])
            for slide_config in slides:
                slide_type = slide_config.get('type', 'content')
                
                if slide_type == 'content':
                    self._add_content_slide(prs, slide_config)
                elif slide_type == 'chart':
                    self._add_chart_slide(prs, slide_config)
                elif slide_type == 'table':
                    self._add_table_slide(prs, slide_config)
                elif slide_type == 'image':
                    self._add_image_slide(prs, slide_config)
                elif slide_type == 'comparison':
                    self._add_comparison_slide(prs, slide_config)
            
            # Save presentation
            prs.save(str(output_path))
            
            log.info(f"Custom PowerPoint presentation created: {output_path}")
            return output_path
            
        except Exception as e:
            log.error(f"Error creating custom PowerPoint presentation: {e}")
            raise

    def create_comprehensive_presentation(self, title: str, author: str = None,
                                        client: str = None, presentation_type: str = "analysis",
                                        sections: List[Dict] = None,
                                        slide_layouts: Dict[str, str] = None,
                                        include_toc: bool = True,
                                        include_agenda: bool = True) -> Dict[str, Any]:
        """
        Create a comprehensive PowerPoint presentation with full slide structure.
        
        Args:
            title: Presentation title
            author: Presentation author
            client: Client name
            presentation_type: Type of presentation
            sections: List of presentation sections
            slide_layouts: Custom slide layouts for different section types
            include_toc: Include table of contents slide
            include_agenda: Include agenda slide
            
        Returns:
            Dictionary containing complete presentation structure
        """
        sections = sections or []
        slide_layouts = slide_layouts or {}
        
        presentation_structure = {
            'title': title,
            'type': 'comprehensive_presentation',
            'metadata': {
                'author': author or 'Siege Analytics',
                'client': client,
                'presentation_type': presentation_type,
                'created_date': datetime.now().isoformat(),
                'version': '1.0'
            },
            'slide_structure': {
                'include_toc': include_toc,
                'include_agenda': include_agenda,
                'slide_layouts': slide_layouts
            },
            'sections': [
                {
                    'type': 'title_slide',
                    'title': title,
                    'subtitle': f"{presentation_type.title()} Presentation",
                    'author': author or 'Siege Analytics',
                    'client': client,
                    'date': datetime.now().strftime('%B %d, %Y'),
                    'level': 0
                }
            ]
        }
        
        # Add table of contents if requested
        if include_toc:
            presentation_structure['sections'].append({
                'type': 'table_of_contents',
                'title': 'Table of Contents',
                'content': self._generate_toc_content(sections),
                'level': 0
            })
        
        # Add agenda if requested
        if include_agenda:
            presentation_structure['sections'].append({
                'type': 'agenda',
                'title': 'Agenda',
                'content': self._generate_agenda_content(sections),
                'level': 0
            })
        
        # Add main sections
        for section in sections:
            section['level'] = section.get('level', 1)
            presentation_structure['sections'].append(section)
        
        return presentation_structure

    def add_slide_section(self, presentation_content: Dict[str, Any], section_type: str,
                          title: str, content: Any, level: int = 1,
                          slide_layout: str = "default",
                          notes: str = None) -> Dict[str, Any]:
        """
        Add a new slide section to an existing presentation.
        
        Args:
            presentation_content: Existing presentation content
            section_type: Type of slide section
            title: Slide title
            content: Slide content
            level: Section level (0=main, 1=section, 2=subsection)
            slide_layout: Custom slide layout
            notes: Speaker notes for the slide
            
        Returns:
            Updated presentation content
        """
        new_section = {
            'type': section_type,
            'title': title,
            'content': content,
            'level': level,
            'slide_layout': slide_layout,
            'notes': notes
        }
        
        if 'sections' not in presentation_content:
            presentation_content['sections'] = []
        
        presentation_content['sections'].append(new_section)
        return presentation_content

    def add_text_slide(self, presentation_content: Dict[str, Any], title: str,
                        text_content: str, level: int = 1,
                        text_style: str = "body", bullet_points: bool = True) -> Dict[str, Any]:
        """
        Add a text slide to the presentation.
        
        Args:
            presentation_content: Existing presentation content
            title: Slide title
            text_content: Text content
            level: Section level
            text_style: Text styling
            bullet_points: Use bullet points for text
            
        Returns:
            Updated presentation content
        """
        text_section = {
            'type': 'text_slide',
            'title': title,
            'content': {
                'text': text_content,
                'style': text_style,
                'bullet_points': bullet_points
            },
            'level': level
        }
        
        return self.add_slide_section(presentation_content, 'text_slide', title, text_section, level)

    def add_chart_slide(self, presentation_content: Dict[str, Any], title: str,
                        charts: List[Any], description: str = "",
                        level: int = 1, layout: str = "single_chart") -> Dict[str, Any]:
        """
        Add a chart slide to the presentation.
        
        Args:
            presentation_content: Existing presentation content
            title: Slide title
            charts: List of chart images
            description: Slide description
            level: Section level
            layout: Chart layout ('single_chart', 'two_charts', 'grid')
            
        Returns:
            Updated presentation content
        """
        chart_section = {
            'type': 'chart_slide',
            'title': title,
            'content': {
                'charts': charts,
                'description': description,
                'layout': layout
            },
            'level': level
        }
        
        return self.add_slide_section(presentation_content, 'chart_slide', title, chart_section, level)

    def add_map_slide(self, presentation_content: Dict[str, Any], title: str,
                      maps: List[Any], map_type: str = "choropleth",
                      description: str = "", level: int = 1) -> Dict[str, Any]:
        """
        Add a map slide to the presentation.
        
        Args:
            presentation_content: Existing presentation content
            title: Slide title
            maps: List of map images
            map_type: Type of maps
            description: Slide description
            level: Section level
            
        Returns:
            Updated presentation content
        """
        map_section = {
            'type': 'map_slide',
            'title': title,
            'content': {
                'maps': maps,
                'map_type': map_type,
                'description': description
            },
            'level': level
        }
        
        return self.add_slide_section(presentation_content, 'map_slide', title, map_section, level)

    def add_table_slide(self, presentation_content: Dict[str, Any], title: str,
                        table_data: Union[List[List], pd.DataFrame],
                        headers: List[str] = None, level: int = 1,
                        table_style: str = "default") -> Dict[str, Any]:
        """
        Add a table slide to the presentation.
        
        Args:
            presentation_content: Existing presentation content
            title: Slide title
            table_data: Table data (list of lists or DataFrame)
            headers: Column headers
            level: Section level
            table_style: Table styling
            
        Returns:
            Updated presentation content
        """
        # Convert DataFrame to list format if needed
        if hasattr(table_data, 'to_dict'):
            if headers is None:
                headers = table_data.columns.tolist()
            table_data = [headers] + table_data.values.tolist()
        elif headers and table_data:
            table_data = [headers] + table_data
        
        table_section = {
            'type': 'table_slide',
            'title': title,
            'content': {
                'data': table_data,
                'headers': headers,
                'style': table_style
            },
            'level': level
        }
        
        return self.add_slide_section(presentation_content, 'table_slide', title, table_section, level)

    def add_comparison_slide(self, presentation_content: Dict[str, Any], title: str,
                            comparison_data: Dict[str, Any], level: int = 1) -> Dict[str, Any]:
        """
        Add a comparison slide to the presentation.
        
        Args:
            presentation_content: Existing presentation content
            title: Slide title
            comparison_data: Data for comparison (before/after, options, etc.)
            level: Section level
            
        Returns:
            Updated presentation content
        """
        comparison_section = {
            'type': 'comparison_slide',
            'title': title,
            'content': comparison_data,
            'level': level
        }
        
        return self.add_slide_section(presentation_content, 'comparison_slide', title, comparison_section, level)

    def add_summary_slide(self, presentation_content: Dict[str, Any], title: str,
                          key_points: List[str], level: int = 1) -> Dict[str, Any]:
        """
        Add a summary slide to the presentation.
        
        Args:
            presentation_content: Existing presentation content
            title: Slide title
            key_points: List of key points to summarize
            level: Section level
            
        Returns:
            Updated presentation content
        """
        summary_section = {
            'type': 'summary_slide',
            'title': title,
            'content': {
                'key_points': key_points
            },
            'level': level
        }
        
        return self.add_slide_section(presentation_content, 'summary_slide', title, summary_section, level)

    def generate_powerpoint_presentation(self, presentation_content: Dict[str, Any],
                                       output_path: str) -> bool:
        """
        Generate a comprehensive PowerPoint presentation.
        
        Args:
            presentation_content: Presentation content structure
            output_path: Output PowerPoint file path
            
        Returns:
            True if successful, False otherwise
        """
        try:
            # Create presentation
            prs = Presentation()
            
            # Apply template if specified
            template_path = presentation_content.get('template_path')
            if template_path and Path(template_path).exists():
                prs = Presentation(template_path)
            
            # Process each section
            for section in presentation_content.get('sections', []):
                slide = self._create_slide_from_section(section, prs)
                if slide:
                    # Add speaker notes if provided
                    if section.get('notes'):
                        notes_slide = slide.notes_slide
                        notes_slide.notes_text_frame.text = section['notes']
            
            # Save presentation
            prs.save(output_path)
            
            log.info(f"PowerPoint presentation generated successfully: {output_path}")
            return True
            
        except Exception as e:
            log.error(f"Error generating PowerPoint presentation: {e}")
            return False

    def _add_title_slide(self, prs: Presentation, title: str):
        """Add a title slide to the presentation."""
        slide_layout = prs.slide_layouts[self.slide_layouts['title']]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Set subtitle
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}\n{self.client_name}"
        
        # Apply formatting
        self._format_title_slide(slide)

    def _add_executive_summary_slide(self, prs: Presentation, summary: str):
        """Add an executive summary slide."""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = "Executive Summary"
        
        # Set content
        content_shape = slide.placeholders[1]
        content_shape.text = summary
        
        # Apply formatting
        self._format_content_slide(slide)

    def _add_metrics_slide(self, prs: Presentation, metrics: Dict[str, Any]):
        """Add a metrics slide with key performance indicators."""
        slide_layout = prs.slide_layouts[self.slide_layouts['two_content']]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = "Key Performance Indicators"
        
        # Add metrics to left content area
        left_content = slide.placeholders[1]
        metrics_text = ""
        
        for metric_name, metric_data in list(metrics.items())[:5]:  # Limit to 5 metrics
            if isinstance(metric_data, dict):
                value = metric_data.get('value', 'N/A')
                status = metric_data.get('status', 'N/A')
                metrics_text += f"• {metric_name}: {value} ({status})\n"
            else:
                metrics_text += f"• {metric_name}: {metric_data}\n"
        
        left_content.text = metrics_text
        
        # Add additional metrics to right content area if available
        if len(metrics) > 5:
            right_content = slide.placeholders[2]
            additional_metrics_text = ""
            
            for metric_name, metric_data in list(metrics.items())[5:]:
                if isinstance(metric_data, dict):
                    value = metric_data.get('value', 'N/A')
                    status = metric_data.get('status', 'N/A')
                    additional_metrics_text += f"• {metric_name}: {value} ({status})\n"
                else:
                    additional_metrics_text += f"• {metric_name}: {metric_data}\n"
            
            right_content.text = additional_metrics_text
        
        # Apply formatting
        self._format_content_slide(slide)

    def _add_charts_slides(self, prs: Presentation, charts: List[Dict[str, Any]]):
        """Add slides with charts."""
        for i, chart_config in enumerate(charts):
            slide_layout = prs.slide_layouts[self.slide_layouts['content']]
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = chart_config.get('title', f'Chart {i+1}')
            
            # Add chart description
            content_shape = slide.placeholders[1]
            content_shape.text = f"Chart type: {chart_config.get('type', 'Unknown')}\n\nData visualization will be added here."
            
            # Apply formatting
            self._format_content_slide(slide)

    def _add_tables_slides(self, prs: Presentation, tables: List[Dict[str, Any]]):
        """Add slides with data tables."""
        for i, table_config in enumerate(tables):
            slide_layout = prs.slide_layouts[self.slide_layouts['content']]
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = table_config.get('title', f'Table {i+1}')
            
            # Add table description
            content_shape = slide.placeholders[1]
            headers = table_config.get('headers', [])
            data = table_config.get('data', [])
            
            content_text = f"Table with {len(headers)} columns and {len(data)} rows\n\n"
            content_text += "Headers: " + ", ".join(headers) + "\n\n"
            content_text += "Data table will be formatted here."
            
            content_shape.text = content_text
            
            # Apply formatting
            self._format_content_slide(slide)

    def _add_insights_slide(self, prs: Presentation, insights: List[str]):
        """Add an insights and recommendations slide."""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = "Key Insights & Recommendations"
        
        # Set content
        content_shape = slide.placeholders[1]
        insights_text = ""
        
        for i, insight in enumerate(insights, 1):
            insights_text += f"{i}. {insight}\n\n"
        
        content_shape.text = insights_text
        
        # Apply formatting
        self._format_content_slide(slide)

    def _add_performance_overview_slide(self, prs: Presentation, overview: str):
        """Add a performance overview slide."""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = "Performance Overview"
        
        # Set content
        content_shape = slide.placeholders[1]
        content_shape.text = overview
        
        # Apply formatting
        self._format_content_slide(slide)

    def _add_metrics_breakdown_slide(self, prs: Presentation, performance_data: Dict[str, Any], metrics: List[str]):
        """Add a metrics breakdown slide."""
        slide_layout = prs.slide_layouts[self.slide_layouts['comparison']]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = "Metrics Breakdown"
        
        # Add metrics to left content area
        left_content = slide.placeholders[1]
        left_metrics_text = ""
        
        for metric in metrics[:len(metrics)//2]:
            if metric in performance_data:
                metric_data = performance_data[metric]
                if isinstance(metric_data, dict):
                    current = metric_data.get('current', 'N/A')
                    target = metric_data.get('target', 'N/A')
                    left_metrics_text += f"• {metric}:\n  Current: {current}\n  Target: {target}\n\n"
                else:
                    left_metrics_text += f"• {metric}: {metric_data}\n\n"
        
        left_content.text = left_metrics_text
        
        # Add remaining metrics to right content area
        right_content = slide.placeholders[2]
        right_metrics_text = ""
        
        for metric in metrics[len(metrics)//2:]:
            if metric in performance_data:
                metric_data = performance_data[metric]
                if isinstance(metric_data, dict):
                    current = metric_data.get('current', 'N/A')
                    target = metric_data.get('target', 'N/A')
                    right_metrics_text += f"• {metric}:\n  Current: {current}\n  Target: {target}\n\n"
                else:
                    right_metrics_text += f"• {metric}: {metric_data}\n\n"
        
        right_content.text = right_metrics_text
        
        # Apply formatting
        self._format_content_slide(slide)

    def _add_trends_slides(self, prs: Presentation, trends: Dict[str, Any]):
        """Add slides with performance trends."""
        for trend_name, trend_data in trends.items():
            slide_layout = prs.slide_layouts[self.slide_layouts['content']]
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = f"{trend_name} Trend"
            
            # Add trend description
            content_shape = slide.placeholders[1]
            content_shape.text = f"Trend analysis for {trend_name}\n\nTrend chart will be added here."
            
            # Apply formatting
            self._format_content_slide(slide)

    def _add_recommendations_slide(self, prs: Presentation, recommendations: List[str]):
        """Add a recommendations slide."""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = "Recommendations"
        
        # Set content
        content_shape = slide.placeholders[1]
        recommendations_text = ""
        
        for i, recommendation in enumerate(recommendations, 1):
            recommendations_text += f"{i}. {recommendation}\n\n"
        
        content_shape.text = recommendations_text
        
        # Apply formatting
        self._format_content_slide(slide)

    def _add_content_slide(self, prs: Presentation, slide_config: Dict[str, Any]):
        """Add a content slide."""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = slide_config.get('title', 'Content Slide')
        
        # Set content
        content_shape = slide.placeholders[1]
        content_shape.text = slide_config.get('content', 'Content will be added here.')
        
        # Apply formatting
        self._format_content_slide(slide)

    def _add_chart_slide(self, prs: Presentation, slide_config: Dict[str, Any]):
        """Add a chart slide."""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = slide_config.get('title', 'Chart Slide')
        
        # Add chart description
        content_shape = slide.placeholders[1]
        caption = slide_config.get('caption', 'Chart will be added here.')
        content_shape.text = caption
        
        # Apply formatting
        self._format_content_slide(slide)

    def _add_table_slide(self, prs: Presentation, slide_config: Dict[str, Any]):
        """Add a table slide."""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = slide_config.get('title', 'Table Slide')
        
        # Add table description
        content_shape = slide.placeholders[1]
        headers = slide_config.get('headers', [])
        data = slide_config.get('data', [])
        
        content_text = f"Table with {len(headers)} columns and {len(data)} rows\n\n"
        content_text += "Headers: " + ", ".join(headers) + "\n\n"
        content_text += "Data table will be formatted here."
        
        content_shape.text = content_text
        
        # Apply formatting
        self._format_content_slide(slide)

    def _add_image_slide(self, prs: Presentation, slide_config: Dict[str, Any]):
        """Add an image slide."""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = slide_config.get('title', 'Image Slide')
        
        # Add image description
        content_shape = slide.placeholders[1]
        caption = slide_config.get('caption', 'Image will be added here.')
        content_shape.text = caption
        
        # Apply formatting
        self._format_content_slide(slide)

    def _add_comparison_slide(self, prs: Presentation, slide_config: Dict[str, Any]):
        """Add a comparison slide."""
        slide_layout = prs.slide_layouts[self.slide_layouts['comparison']]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = slide_config.get('title', 'Comparison Slide')
        
        # Add left content
        left_content = slide.placeholders[1]
        left_content.text = slide_config.get('left_content', 'Left content')
        
        # Add right content
        right_content = slide.placeholders[2]
        right_content.text = slide_config.get('right_content', 'Right content')
        
        # Apply formatting
        self._format_content_slide(slide)

    def _format_title_slide(self, slide):
        """Apply formatting to title slide."""
        try:
            # Format title
            title_shape = slide.shapes.title
            title_shape.text_frame.paragraphs[0].font.size = Pt(44)
            title_shape.text_frame.paragraphs[0].font.bold = True
            title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Format subtitle
            subtitle_shape = slide.placeholders[1]
            for paragraph in subtitle_shape.text_frame.paragraphs:
                paragraph.font.size = Pt(24)
                paragraph.alignment = PP_ALIGN.CENTER
                
        except Exception as e:
            log.warning(f"Could not apply formatting to title slide: {e}")

    def _format_content_slide(self, slide):
        """Apply formatting to content slide."""
        try:
            # Format title
            title_shape = slide.shapes.title
            title_shape.text_frame.paragraphs[0].font.size = Pt(36)
            title_shape.text_frame.paragraphs[0].font.bold = True
            
            # Format content
            content_shape = slide.placeholders[1]
            for paragraph in content_shape.text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                
        except Exception as e:
            log.warning(f"Could not apply formatting to content slide: {e}")

    def create_presentation_from_dataframe(self, df: pd.DataFrame,
                                         presentation_title: str = "",
                                         max_slides: int = 10) -> Path:
        """
        Create a PowerPoint presentation from a pandas DataFrame.
        
        Args:
            df: Pandas DataFrame
            presentation_title: Title of the presentation
            max_slides: Maximum number of slides to create
            
        Returns:
            Path to the generated PowerPoint file
        """
        try:
            # Generate filename
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"{self.client_name.lower().replace(' ', '_')}_dataframe_presentation_{timestamp}.pptx"
            output_path = self.output_dir / filename
            
            # Create presentation
            prs = Presentation()
            
            # Add title slide
            title = presentation_title or f"{self.client_name} Data Analysis"
            self._add_title_slide(prs, title)
            
            # Add data overview slide
            self._add_data_overview_slide(prs, df)
            
            # Add summary statistics slide
            self._add_summary_stats_slide(prs, df)
            
            # Add sample data slides
            self._add_sample_data_slides(prs, df, max_slides - 4)
            
            # Save presentation
            prs.save(str(output_path))
            
            log.info(f"DataFrame PowerPoint presentation created: {output_path}")
            return output_path
            
        except Exception as e:
            log.error(f"Error creating DataFrame PowerPoint presentation: {e}")
            raise

    def _add_data_overview_slide(self, prs: Presentation, df: pd.DataFrame):
        """Add a data overview slide."""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = "Data Overview"
        
        # Set content
        content_shape = slide.placeholders[1]
        content_text = f"Dataset Information:\n\n"
        content_text += f"• Number of rows: {len(df):,}\n"
        content_text += f"• Number of columns: {len(df.columns)}\n"
        content_text += f"• Data types: {', '.join(df.dtypes.astype(str).unique())}\n"
        content_text += f"• Memory usage: {df.memory_usage(deep=True).sum() / 1024 / 1024:.2f} MB"
        
        content_shape.text = content_text
        
        # Apply formatting
        self._format_content_slide(slide)

    def _add_summary_stats_slide(self, prs: Presentation, df: pd.DataFrame):
        """Add a summary statistics slide."""
        slide_layout = prs.slide_layouts[self.slide_layouts['two_content']]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = "Summary Statistics"
        
        # Get numeric columns
        numeric_cols = df.select_dtypes(include=['number']).columns
        
        # Add left content with first half of statistics
        left_content = slide.placeholders[1]
        left_text = ""
        
        for col in numeric_cols[:len(numeric_cols)//2]:
            stats = df[col].describe()
            left_text += f"{col}:\n"
            left_text += f"  Mean: {stats['mean']:.2f}\n"
            left_text += f"  Std: {stats['std']:.2f}\n"
            left_text += f"  Min: {stats['min']:.2f}\n"
            left_text += f"  Max: {stats['max']:.2f}\n\n"
        
        left_content.text = left_text
        
        # Add right content with second half of statistics
        if len(numeric_cols) > 1:
            right_content = slide.placeholders[2]
            right_text = ""
            
            for col in numeric_cols[len(numeric_cols)//2:]:
                stats = df[col].describe()
                right_text += f"{col}:\n"
                right_text += f"  Mean: {stats['mean']:.2f}\n"
                right_text += f"  Std: {stats['std']:.2f}\n"
                right_text += f"  Min: {stats['min']:.2f}\n"
                right_text += f"  Max: {stats['max']:.2f}\n\n"
            
            right_content.text = right_text
        
        # Apply formatting
        self._format_content_slide(slide)

    def _add_sample_data_slides(self, prs: Presentation, df: pd.DataFrame, max_slides: int):
        """Add slides with sample data."""
        if df.empty or max_slides <= 0:
            return
        
        # Calculate rows per slide
        rows_per_slide = max(1, min(20, len(df) // max_slides))
        
        for i in range(0, min(len(df), max_slides * rows_per_slide), rows_per_slide):
            slide_layout = prs.slide_layouts[self.slide_layouts['content']]
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = f"Sample Data (Rows {i+1}-{min(i+rows_per_slide, len(df))})"
            
            # Set content
            content_shape = slide.placeholders[1]
            
            # Get sample data
            sample_df = df.iloc[i:i+rows_per_slide]
            
            # Create formatted table text
            content_text = "Columns: " + " | ".join(df.columns.tolist()) + "\n\n"
            
            for idx, row in sample_df.iterrows():
                row_text = " | ".join([str(val)[:20] for val in row.values])
                content_text += f"{row_text}\n"
            
            content_shape.text = content_text
            
            # Apply formatting
            self._format_content_slide(slide)

    def _create_slide_from_section(self, section: Dict[str, Any], prs: Presentation) -> Any:
        """
        Create a PowerPoint slide from a section definition.
        
        Args:
            section: Section dictionary
            prs: PowerPoint presentation object
            
        Returns:
            Created slide object
        """
        section_type = section.get('type', 'text_slide')
        
        try:
            if section_type == 'title_slide':
                return self._create_title_slide(section, prs)
            elif section_type == 'table_of_contents':
                return self._create_toc_slide(section, prs)
            elif section_type == 'agenda':
                return self._create_agenda_slide(section, prs)
            elif section_type == 'text_slide':
                return self._create_text_slide(section, prs)
            elif section_type == 'chart_slide':
                return self._create_chart_slide(section, prs)
            elif section_type == 'map_slide':
                return self._create_map_slide(section, prs)
            elif section_type == 'table_slide':
                return self._create_table_slide(section, prs)
            elif section_type == 'comparison_slide':
                return self._create_comparison_slide(section, prs)
            elif section_type == 'summary_slide':
                return self._create_summary_slide(section, prs)
            else:
                # Default to text slide
                return self._create_text_slide(section, prs)
                
        except Exception as e:
            log.error(f"Error creating slide for section {section_type}: {e}")
            return None

    def _create_title_slide(self, section: Dict[str, Any], prs: Presentation) -> Any:
        """Create a title slide."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = section.get('title', '')
        
        # Set subtitle
        subtitle = slide.placeholders[1]
        subtitle.text = section.get('subtitle', '')
        
        # Add additional information
        if 'author' in section or 'client' in section or 'date' in section:
            # Add text box for additional info
            left = Inches(1)
            top = Inches(6)
            width = Inches(8)
            height = Inches(2)
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            
            if 'author' in section:
                p = tf.add_paragraph()
                p.text = f"Author: {section['author']}"
                p.font.size = Pt(14)
            
            if 'client' in section:
                p = tf.add_paragraph()
                p.text = f"Client: {section['client']}"
                p.font.size = Pt(14)
            
            if 'date' in section:
                p = tf.add_paragraph()
                p.text = f"Date: {section['date']}"
                p.font.size = Pt(14)
        
        return slide

    def _create_toc_slide(self, section: Dict[str, Any], prs: Presentation) -> Any:
        """Create a table of contents slide."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = section.get('title', 'Table of Contents')
        
        # Add content
        content = slide.placeholders[1]
        tf = content.text_frame
        
        # Add TOC items
        toc_content = section.get('content', [])
        for item in toc_content:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.level = 0
        
        return slide

    def _create_agenda_slide(self, section: Dict[str, Any], prs: Presentation) -> Any:
        """Create an agenda slide."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = section.get('title', 'Agenda')
        
        # Add content
        content = slide.placeholders[1]
        tf = content.text_frame
        
        # Add agenda items
        agenda_content = section.get('content', [])
        for item in agenda_content:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.level = 0
        
        return slide

    def _create_text_slide(self, section: Dict[str, Any], prs: Presentation) -> Any:
        """Create a text slide."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = section.get('title', '')
        
        # Add content
        content = slide.placeholders[1]
        tf = content.text_frame
        
        text_content = section.get('content', {}).get('text', '')
        use_bullets = section.get('content', {}).get('bullet_points', True)
        
        if use_bullets and '\n' in text_content:
            # Split into bullet points
            lines = text_content.split('\n')
            for i, line in enumerate(lines):
                if line.strip():
                    p = tf.add_paragraph()
                    p.text = line.strip()
                    p.font.size = Pt(18)
                    p.level = 0
        else:
            # Single paragraph
            p = tf.add_paragraph()
            p.text = text_content
            p.font.size = Pt(18)
        
        return slide

    def _create_chart_slide(self, section: Dict[str, Any], prs: Presentation) -> Any:
        """Create a chart slide."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = section.get('title', '')
        
        # Add charts
        charts = section.get('content', {}).get('charts', [])
        layout = section.get('content', {}).get('layout', 'single_chart')
        
        if layout == 'single_chart' and charts:
            # Single chart
            chart = charts[0]
            self._add_chart_to_slide(chart, slide, Inches(1), Inches(2), Inches(8), Inches(5))
        elif layout == 'two_charts' and len(charts) >= 2:
            # Two charts side by side
            self._add_chart_to_slide(charts[0], slide, Inches(0.5), Inches(2), Inches(4), Inches(4))
            self._add_chart_to_slide(charts[1], slide, Inches(5.5), Inches(2), Inches(4), Inches(4))
        elif layout == 'grid' and len(charts) >= 4:
            # 2x2 grid
            positions = [
                (Inches(0.5), Inches(2), Inches(3.5), Inches(2.5)),
                (Inches(4.5), Inches(2), Inches(3.5), Inches(2.5)),
                (Inches(0.5), Inches(4.8), Inches(3.5), Inches(2.5)),
                (Inches(4.5), Inches(4.8), Inches(3.5), Inches(2.5))
            ]
            for i, chart in enumerate(charts[:4]):
                left, top, width, height = positions[i]
                self._add_chart_to_slide(chart, slide, left, top, width, height)
        
        return slide

    def _create_map_slide(self, section: Dict[str, Any], prs: Presentation) -> Any:
        """Create a map slide."""
        # Similar to chart slide but optimized for maps
        return self._create_chart_slide(section, prs)

    def _create_table_slide(self, section: Dict[str, Any], prs: Presentation) -> Any:
        """Create a table slide."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = section.get('title', '')
        
        # Add table
        table_data = section.get('content', {}).get('data', [])
        if table_data:
            rows = len(table_data)
            cols = len(table_data[0]) if table_data else 1
            
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(4)
            
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Populate table
            for i, row in enumerate(table_data):
                for j, cell_value in enumerate(row):
                    table.cell(i, j).text = str(cell_value)
        
        return slide

    def _create_comparison_slide(self, section: Dict[str, Any], prs: Presentation) -> Any:
        """Create a comparison slide."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = section.get('title', '')
        
        # Add comparison content
        comparison_data = section.get('content', {})
        
        # Create two-column layout
        left = Inches(0.5)
        top = Inches(2)
        width = Inches(4)
        height = Inches(4)
        
        # Left column
        left_box = slide.shapes.add_textbox(left, top, width, height)
        left_tf = left_box.text_frame
        left_tf.text = comparison_data.get('left_title', 'Before')
        
        # Right column
        right_box = slide.shapes.add_textbox(Inches(5.5), top, width, height)
        right_tf = right_box.text_frame
        right_tf.text = comparison_data.get('right_title', 'After')
        
        return slide

    def _create_summary_slide(self, section: Dict[str, Any], prs: Presentation) -> Any:
        """Create a summary slide."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = section.get('title', 'Summary')
        
        # Add key points
        content = slide.placeholders[1]
        tf = content.text_frame
        
        key_points = section.get('content', {}).get('key_points', [])
        for point in key_points:
            p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(18)
            p.level = 0
        
        return slide

    def _generate_toc_content(self, sections: List[Dict]) -> List[str]:
        """Generate table of contents content from sections."""
        toc_items = []
        for section in sections:
            if section.get('type') not in ['title_slide', 'table_of_contents', 'agenda']:
                toc_items.append(section.get('title', 'Untitled'))
        return toc_items

    def _generate_agenda_content(self, sections: List[Dict]) -> List[str]:
        """Generate agenda content from sections."""
        agenda_items = []
        for section in sections:
            if section.get('type') not in ['title_slide', 'table_of_contents', 'agenda']:
                agenda_items.append(section.get('title', 'Untitled'))
        return agenda_items
